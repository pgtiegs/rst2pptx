import unittest
import glob

import docutils.core
import docutils.io

from lxml import etree
import pptx
import rst2pptx

class Test_rst2pptx(unittest.TestCase):
    def test_basic_read_write(self):
        sample = """
=====
Title
=====

My Title Slide
==============
* This is slide

Second Slide
============

* Hello

sub-section
-----------
My Subsection

"""
        writer = rst2pptx.PowerPointWriter()
        writer.document = docutils.core.publish_doctree(sample)
        writer.presentation = pptx.Presentation()
        writer.translate()

        self.assertIsInstance(writer.presentation, pptx.presentation.Presentation )

        self.assertEqual(writer.presentation.slides[0].slide_layout.name, "Title Slide")
        
       
        print([x.name for x in writer.presentation.slide_layouts])
        writer.presentation.save("test.pptx")

    def base_slides(self, feature):
        writer = rst2pptx.PowerPointWriter()
        writer.presentation = pptx.Presentation()
        with open("test/{}.rst".format(feature), 'r') as fd:
            writer.document = docutils.core.publish_doctree(fd.read())
        writer.translate()
        base_slides = [etree.parse(x).getroot() for x in glob.glob("test/{}/*.xml".format(feature))] 
        for base, slide in zip(base_slides, writer.presentation.slides):
            print(slide.element)
            self.assertEqual(etree.tostring(base), etree.tostring(slide.element))

    def test_bullets(self):
        self.base_slides("bullets")

    def test_enumerated_lists(self):
        self.base_slides("enumerated_list")

    def test_image_from_uri(self):
        self.base_slides("image_from_uri")

    def test_text_in_title_slide(self):
        self.base_slides("text_in_title_slide")

    def test_hyperlinks(self):
        self.base_slides("hyperlink")

    def test_subsections(self):
        self.base_slides("subsections")

    def test_definitionlist(self):
        self.base_slides("definitionlist")

    def test_classes(self):
        self.base_slides("classes")

if __name__ == '__main__':
    unittest.main()
