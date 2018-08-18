#!/usr/bin/env python
# encoding: utf-8

# Copyright (C) 2016 Steven Myint
#
# Permission is hereby granted, free of charge, to any person obtaining
# a copy of this software and associated documentation files (the
# "Software"), to deal in the Software without restriction, including
# without limitation the rights to use, copy, modify, merge, publish,
# distribute, sublicense, and/or sell copies of the Software, and to
# permit persons to whom the Software is furnished to do so, subject to
# the following conditions:
#
# The above copyright notice and this permission notice shall be
# included in all copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND,
# EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF
# MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND
# NONINFRINGEMENT. IN NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS
# BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER LIABILITY, WHETHER IN AN
# ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM, OUT OF OR IN
# CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

"""Converts reStructuredText to PowerPoint."""

import io
import os
import sys
import urllib
import logging

import docutils.core
import docutils.nodes
import docutils.utils
import pptx

from lxml import etree
__version__ = '0.3'

logging.basicConfig(level=logging.WARNING)

TITLE_BUFFER = pptx.util.Inches(2.)
MARGIN = pptx.util.Inches(1.)

def setBuNone(paragraph):
    etree.SubElement(paragraph._pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")

def setBuAutoNum(paragraph):
    paragraph._pPr.attrib['marL'] = "427789"
    paragraph._pPr.attrib['indent'] = "-427789"
    e = etree.SubElement(paragraph._pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
    e.attrib["type"] = "arabicPeriod"
    e.attrib["startAt"] = "1"

class PowerPointTranslator(docutils.nodes.NodeVisitor):

    """A translator for converting docutils elements to PowerPoint."""

    def __init__(self, document, presentation):
        docutils.nodes.NodeVisitor.__init__(self, document)

        self.bullet_level = 0
        self.bullet_list = False
        self.enum_list = False
        self.presentation = presentation
        self.slides = self.presentation.slides
        self.table_rows = None
        self.title_slide = True
        self.section_level = 0

    def visit_document(self, node):
        pass

    def depart_document(self, node):
        pass

    def visit_decoration(self,node):
        logging.debug("-> decoration")
    
    def depart_decoration(self,node):
        logging.debug("decoration ->")

    def visit_footer(self,node):
        logging.debug("-> footer")
    
    def depart_footer(self,node):
        logging.debug("footer ->")

    def visit_substitution_definition(self,node):
        logging.debug("-> substitution_definition")
        raise docutils.nodes.SkipNode
    
    def depart_substitution_definition(self,node):
        logging.debug("substitution_definition ->")

    def visit_author(self,node):
        logging.debug("-> author")
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
    
    def depart_author(self,node):
        logging.debug("author ->")

    def visit_date(self,node):
        logging.debug("-> date")
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
    
    def depart_date(self,node):
        logging.debug("date ->")

    def visit_version(self,node):
        logging.debug("-> version")
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
    
    def depart_version(self,node):
        logging.debug("version ->")

    def visit_status(self,node):
        logging.debug("-> status")
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = "Status: "
        
    
    def depart_status(self,node):
        logging.debug("status ->")

    def visit_copyright(self,node):
        logging.debug("-> copyright")
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.add_paragraph()
    
    def depart_copyright(self,node):
        logging.debug("copyright ->")

    def visit_docinfo(self,node):
        logging.debug("-> docinfo")
    
    def depart_docinfo(self,node):
        logging.debug("docinfo ->")

    def visit_docinfo_item(self, node, name):
        pass

    def visit_image(self, node):
        uri = node.attributes['uri']
        if '://' in uri:
            if sys.version_info[0] < 3:
                self.document.reporter.warning(
                    'Downloading images requires Python 3 or greater')
                return

            try:
                with urllib.request.urlopen(uri) as input_file:
                    image_file = io.BytesIO(input_file.read())
            except urllib.error.HTTPError as e:
                self.document.reporter.warning(
                    'Could not access {}'.format(uri))
                return
        else:
            document_filename = docutils.utils.get_source_line(node)[0]
            if document_filename and document_filename != '<stdin>':
                root_path = os.path.dirname(document_filename)
            else:
                root_path = os.getcwd()
            image_file = os.path.join(root_path, uri)
            uri = image_file

        try:
            picture = self.slides[-1].shapes.add_picture(
                image_file,
                left=0,
                top=0)
        except IOError:
            self.document.reporter.warning(
                'Could not load image {}'.format(uri))
            return

        center_picture(picture, self.presentation)

    def depart_image(self, node):
        pass

    def visit_figure(self,node):
        logging.debug("-> figure")
    
    def depart_figure(self,node):
        logging.debug("figure ->")

    def visit_caption(self,node):
        logging.debug("-> caption")
    
    def depart_caption(self,node):
        logging.debug("caption ->")

    def visit_Text(self, node):
        logging.debug("visiting text")

        logging.debug("text parent = {}".format(node.parent.tagname))
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.paragraphs[-1]
        run = paragraph.add_run()
        run.text = node.astext()

    def depart_Text(self, node):
        logging.debug("departing text")
        pass

    def visit_list_item(self, node):
        logging.debug("visiting list_item")
        
    def depart_list_item(self, node):
        logging.debug("departing list_item")
        pass

    def visit_paragraph(self, node):
        logging.debug("visiting paragraph")

        shapes = self.slides[-1].shapes

        if self.title_slide and not shapes[-1].text:
            # This must be the empty text box for the subtitle.
            pass
        else:
            text_frame = self.slides[-1].shapes.placeholders[1].text_frame
            paragraph = text_frame.add_paragraph()
            if not self.bullet_list:
                if self.enum_list:
                    setBuAutoNum(paragraph)
                else:
                    setBuNone(paragraph)
            if self.bullet_list:
                paragraph.level = self.bullet_level

    def depart_paragraph(self, node):
        logging.debug("departing paragraph")
        pass

    def visit_section(self, node):
        logging.debug("-> section")
        logging.debug(self.section_level)
        if self.section_level == 0:
            self.title_slide = False
            self.slides.add_slide(self.presentation.slide_layouts[1])
        else:
            logging.debug("SubSection")

        self.section_level += 1

    def depart_section(self, node):
        self.section_level -= 1
        logging.debug("section ->")

    def visit_title(self, node):
        logging.debug("visiting title")
        logging.debug("text parent = {}".format(type(node.parent.tagname)))
        if len(self.slides):
            logging.debug("{} {}".format(node.astext(), self.section_level))
            #self.slides[-1].shapes.title.text = node.astext()
            if self.section_level == 1:
                self.slides[-1].shapes.title.text = node.astext()
            elif self.section_level >= 1:
                text_frame = self.slides[-1].shapes.placeholders[1].text_frame
                paragraph = text_frame.add_paragraph()
                setBuNone(paragraph)
                run = paragraph.add_run()
                run.text = node.astext()
                run.font.bold = True

            elif node.parent.tagname == "topic":
                logging.debug("in topic")
                text_frame = self.slides[-1].shapes.placeholders[1].text_frame
                paragraph = text_frame.add_paragraph()
                run = paragraph.add_run()
                run.text = node.astext()
                run.font.bold = True
        else:
            # Title slide.
            slide = self.slides.add_slide(self.presentation.slide_layouts[0])
            slide.shapes.title.text = node.astext()
            self.title_slide = True
            # TODO: Author.
        raise docutils.nodes.SkipNode

    def depart_title(self, node):
        logging.debug("departing title")
        pass

    def visit_literal_block(self, node):
        pass

    def depart_literal_block(self, node):
        pass

    def visit_literal(self,node):
        logging.debug("-> literal")

    def depart_literal(self,node):
        logging.debug("literal ->")

    def visit_definition_list(self,node):
        logging.debug("-> definition_list")

    def depart_definition_list(self,node):
        logging.debug("definition_list->")

    def visit_definition_list_item(self,node):
        logging.debug("-> definition_list_item")

    def depart_definition_list_item(self,node):
        logging.debug("definition_list_item ->")

    def visit_term(self,node):
        logging.debug("-> term")
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.add_paragraph()
        setBuNone(paragraph)
        logging.debug("term: {}".format(node.astext()))

    def depart_term(self,node):
        logging.debug("term ->")

    def visit_definition(self,node):
        logging.debug("-> definition")
        #text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        #paragraph = text_frame.add_paragraph()

    def depart_definition(self,node):
        logging.debug("definition ->")
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.paragraphs[-1]
        level = paragraph.level
        paragraph.level = level +1

    def visit_block_quote(self,node):
        logging.debug("-> block_quote")

    def depart_block_quote(self,node):
        logging.debug("block_quote ->")

    def visit_inline(self,node):
        logging.debug("-> inline")

    def depart_inline(self,node):
        logging.debug("inline ->")

    def visit_topic(self,node):
        logging.debug("-> topic")
        logging.debug(node)

    def depart_topic(self,node):
        logging.debug("topic ->")

    def visit_transition(self,node):
        logging.debug("-> transition")

    def depart_transition(self,node):
        logging.debug("transition ->")

    def visit_bullet_list(self, node):
        if self.bullet_list:
            self.bullet_level += 1
        else:
            self.bullet_list = True
        logging.debug("visiting bullet_level {}".format(self.bullet_level))



    def depart_bullet_list(self, node):
        if self.bullet_level == 0:
            self.bullet_list = False
        if self.bullet_list:
            self.bullet_level -= 1
        logging.debug("departing bullet_level {}".format(self.bullet_level))
        assert self.bullet_level >= 0 

    def visit_enumerated_list(self, node):
        logging.debug("-> enumerated_list")
        if self.enum_list:
            self.bullet_level += 1
        else:
            self.enum_list = True
        logging.debug("visiting bullet_level {}".format(self.bullet_level))

    def depart_enumerated_list(self, node):
        logging.debug("enumerated_list ->")
        if self.bullet_level == 0:
            self.enum_list = False
        if self.enum_list:
            self.bullet_level -= 1
        logging.debug("departing bullet_level {}".format(self.bullet_level))
        assert self.bullet_level >= 0 

    def visit_tgroup(self, node):
        self.table_rows = []

    def depart_tgroup(self, node):
        if self.table_rows and self.table_rows[0]:
            table = self.slides[-1].shapes.add_table(
                rows=len(self.table_rows),
                cols=len(self.table_rows[0]),
                left=MARGIN,
                top=TITLE_BUFFER,
                width=self.presentation.slide_width - 2 * MARGIN,
                height=self.presentation.slide_height - 2 * TITLE_BUFFER).table

            for (row_index, row) in enumerate(self.table_rows):
                for (col_index, col) in enumerate(row):
                    table.cell(row_idx=row_index, col_idx=col_index).text = col

            self.table_rows = None

    def visit_tbody(self,node):
        logging.debug("-> tbody")
    
    def depart_tbody(self,node):
        logging.debug("tbody ->")

    def visit_thead(self,node):
        logging.debug("-> thead")
    
    def depart_thead(self,node):
        logging.debug("thead ->")

    def visit_table(self,node):
        logging.debug("-> table")
    
    def depart_table(self,node):
        logging.debug("table ->")

    def visit_colspec(self,node):
        logging.debug("-> colspec")
    
    def depart_colspec(self,node):
        logging.debug("colspec ->")

    def visit_row(self, node):
        assert self.table_rows is not None
        self.table_rows.append([])

    def depart_row(self, node):
        pass

    def visit_entry(self, node):
        self.table_rows[-1].append(node.astext())
        raise docutils.nodes.SkipNode

    def depart_entry(self, node):
        pass

    def visit_reference(self, node):
        logging.debug("visiting reference")

    def depart_reference(self, node):    
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.paragraphs[-1]
        run = paragraph.runs[-1]
    
        run.hyperlink.address = node.attributes.get('refuri')
        logging.debug("departing reference")

    def visit_strong(self, node):
        logging.debug("visiting strong")

    def depart_strong(self, node):
        text_frame = self.slides[-1].shapes.placeholders[1].text_frame
        paragraph = text_frame.paragraphs[-1]
        run = paragraph.runs[-1]

        run.font.bold = True
        logging.debug("departing strong")

    def visit_target(self, node):
        pass

    def depart_target(self, node):
        pass

    def unknown_visit(self, node):
        self.document.reporter.warning('unknown_visit({})'.format(node))

    def unknown_departure(self, node):
        self.document.reporter.warning('unknown_departure({})'.format(node))

    def astext(self):
        pass


def center_picture(picture, presentation):
    picture.left = (presentation.slide_width - picture.width) // 2

    slide_height = presentation.slide_height - TITLE_BUFFER
    picture.top = (slide_height - picture.height) // 2 + TITLE_BUFFER


class PowerPointWriter(docutils.core.writers.Writer):

    """A docutils writer that produces PowerPoint."""

    settings_spec = (
        'PowerPoint options',
        None,
        (
            (
                'PowerPoint template.',
                ['--pptx-template'],
                {'default': None}
            ),
        )
    )

    def __init__(self):
        docutils.core.writers.Writer.__init__(self)

        self.presentation = None
        self.translator_class = PowerPointTranslator

    def translate(self):
        assert self.presentation
        visitor = self.translator_class(document=self.document,
                                        presentation=self.presentation)
        self.document.walkabout(visitor)

    def write(self, document, destination):
        self.document = document
        self.presentation = pptx.Presentation(document.settings.pptx_template)

        self.language = docutils.languages.get_language(
            document.settings.language_code,
            document.reporter)

        self.translate()

        if destination.destination is None:
            self.presentation.save(destination.destination_path)
        else:
            stream = io.BytesIO()
            self.presentation.save(stream)
            destination.write(stream.getvalue())


def main():
    docutils.core.publish_cmdline(
        writer=PowerPointWriter(),
        description='Generates PowerPoint presentations.  ' +
                    docutils.core.default_description,
        settings_overrides={'halt_level': docutils.utils.Reporter.ERROR_LEVEL})


if __name__ == '__main__':
    main()
